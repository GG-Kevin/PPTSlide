#!/usr/bin/env python3
"""회사 표준 포맷(templates/company_format.pptx)에서 포맷 스펙을 추출한다.

산출물: format/format_spec.json
포맷 가디언 에이전트가 이 스크립트를 돌려 스펙을 갱신하고,
빌더/검증 에이전트는 이 JSON만 신뢰한다.
"""
import json
import re
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "company_format.pptx"
OUT = ROOT / "format" / "format_spec.json"

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def emu_to_in(v):
    return round(v / 914400, 4)


def extract():
    from pptx import Presentation

    prs = Presentation(str(TEMPLATE))
    spec = {
        "template": str(TEMPLATE.relative_to(ROOT)),
        "slide_size_emu": {"cx": prs.slide_width, "cy": prs.slide_height},
        "slide_size_in": {
            "w": emu_to_in(prs.slide_width),
            "h": emu_to_in(prs.slide_height),
        },
        "layouts": {},
    }

    for layout in prs.slide_masters[0].slide_layouts:
        shapes = []
        for sh in layout.shapes:
            ph = sh.element.find(f".//{P}ph")
            shapes.append(
                {
                    "name": sh.name,
                    "shape_type": str(sh.shape_type),
                    "ph": dict(ph.attrib) if ph is not None else None,
                    "left": sh.left,
                    "top": sh.top,
                    "width": sh.width,
                    "height": sh.height,
                }
            )
        spec["layouts"][layout.name] = {"shapes": shapes}

    # 본문 레이아웃의 문단 수준별 스타일(파라그래프 레벨 -> 폰트/크기/불릿)
    with zipfile.ZipFile(TEMPLATE) as z:
        body_layout = z.read("ppt/slideLayouts/slideLayout1.xml").decode("utf-8")
        theme = z.read("ppt/theme/theme1.xml").decode("utf-8")
        motif = z.read("ppt/media/image1.png")

    levels = {}
    body_block = body_layout.split('name="텍스트 개체 틀 8"')[1]
    body_block = body_block.split("<a:lstStyle>")[1].split("</a:lstStyle>")[0]
    for m in re.finditer(r"<a:lvl(\d)pPr([^>]*)>(.*?)</a:lvl\1pPr>", body_block, re.S):
        lvl, attrs, inner = m.group(1), m.group(2), m.group(3)
        sz = re.search(r'<a:defRPr[^>]*\bsz="(\d+)"', inner)
        b = re.search(r'<a:defRPr[^>]*\bb="(\d)"', inner)
        latin = re.search(r'<a:latin typeface="([^"]+)"', inner)
        buchar = re.search(r'<a:buChar char="([^"]+)"', inner)
        bufont = re.search(r'<a:buFont typeface="([^"]+)"', inner)
        lnspc = re.search(r'<a:lnSpc><a:spcPts val="(\d+)"/>', inner)
        marl = re.search(r'marL="(-?\d+)"', attrs)
        ind = re.search(r'indent="(-?\d+)"', attrs)
        levels[f"lvl{lvl}"] = {
            "sz": int(sz.group(1)) if sz else None,
            "bold": b.group(1) == "1" if b else None,
            "font": latin.group(1) if latin else None,
            "bullet_char": buchar.group(1) if buchar else None,
            "bullet_font": bufont.group(1) if bufont else None,
            "line_spacing_pts": int(lnspc.group(1)) / 100 if lnspc else None,
            "marL": int(marl.group(1)) if marl else None,
            "indent": int(ind.group(1)) if ind else None,
        }

    title_block = body_layout.split('name="제목 1"')[1]
    title_block = title_block.split("<a:lstStyle>")[1].split("</a:lstStyle>")[0]
    spec["title_style"] = {
        "sz": int(re.search(r'\bsz="(\d+)"', title_block).group(1)),
        "bold": bool(re.search(r'\bb="1"', title_block)),
        "font": re.search(r'<a:latin typeface="([^"]+)"', title_block).group(1),
        "align": "l",
    }
    spec["body_levels"] = levels

    # 테마 색
    clr = {}
    scheme = re.search(r"<a:clrScheme.*?</a:clrScheme>", theme, re.S).group(0)
    for m in re.finditer(r"<a:(\w+)><a:(?:srgbClr val|sysClr val)=\"(\w+)\"", scheme):
        clr[m.group(1)] = m.group(2)
    spec["theme_colors"] = clr
    spec["theme_fonts"] = {
        "major_latin": re.search(r'<a:majorFont><a:latin typeface="([^"]*)"', theme).group(1),
        "minor_latin": re.search(r'<a:minorFont><a:latin typeface="([^"]*)"', theme).group(1),
    }

    # 상단 모티프 바(가로모티프색.png)의 색 구성 — 검증 에이전트의 기준값
    from PIL import Image
    import io

    im = Image.open(io.BytesIO(motif)).convert("RGBA")
    w, _h = im.size
    row = [im.getpixel((x, im.size[1] // 2)) for x in range(w)]
    opaque = [p for p in row if p[3] > 200]
    navy = opaque[0][:3]
    green = opaque[-1][:3]
    split = next(i for i, p in enumerate(row) if p[3] > 200 and p[:3] == green) / w
    spec["motif_bar"] = {
        "image": "ppt/media/image1.png",
        "navy_rgb": list(navy),
        "green_rgb": list(green),
        "navy_hex": "%02X%02X%02X" % navy,
        "green_hex": "%02X%02X%02X" % green,
        "green_starts_at_ratio": round(split, 4),
        "top_emu": 549276,
        "height_emu": 47625,
    }

    spec["body_slide_layout"] = "3. 본문 및 내용"
    spec["rules"] = [
        "새 슬라이드는 반드시 templates/company_format.pptx 를 열어 그 안의 레이아웃으로 추가한다. 새 Presentation()을 만들지 않는다.",
        "본문 슬라이드는 레이아웃 '3. 본문 및 내용'을 쓴다.",
        "제목은 title 플레이스홀더에, 대분류 항목은 body(idx=10) 플레이스홀더 lvl1(❖)에 넣는다.",
        "상단 모티프 바와 우하단 페이지 번호는 레이아웃이 제공한다. 슬라이드에 다시 그리지 않는다.",
        "본문 글꼴은 '현대하모니 L'을 유지한다. 임의의 글꼴/색을 도입하지 않는다.",
        "본문 글자색은 테마 tx1(검정)을 유지한다.",
    ]
    return spec


if __name__ == "__main__":
    spec = extract()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"wrote {OUT}")
    json.dump(spec.get("body_levels"), sys.stdout, ensure_ascii=False, indent=2)
    print()
