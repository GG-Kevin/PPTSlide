#!/usr/bin/env python3
"""블록 조립 방식 슬라이드 렌더러.

초안(draft)을 블록 목록으로 옮기면 그대로 슬라이드가 된다.
블록은 각자 `box`(인치, 슬라이드 절대좌표)를 가지며 서로 겹치지 않게 배치한다.

지원 블록
  image  : 이미지 1장 (종횡비 보존 · 프레임 · 캡션)   ← 프로그램 스냅샷용
  images : 이미지 여러 장 가로 배치
  table  : 실제 PowerPoint 표 (파워포인트에서 그대로 편집 가능)
  text   : 수준별 항목 텍스트
  note   : 핵심 메시지 바
"""
from pathlib import Path

from PIL import Image as PILImage
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Pt

import shapes as S

IN = 914400

INK = "1B2A41"
BODY = "44536A"
MUTED = "8A94A6"
HAIRLINE = "E3E9F0"
CARD_BG = "F7F9FC"
CARD_LN = "E4EAF2"
BRAND_NAVY = "014099"
TABLE_HEAD = "2C3E5D"
FONT = "맑은 고딕"

MIN_DPI = 110   # 스냅샷이 흐려 보이기 시작하는 경계
MIN_ROW_H = int(0.22 * IN)   # 9.5pt 본문이 들어가는 표 행의 실질 최소 높이


def emu_box(box):
    return tuple(int(round(v * IN)) for v in box)


# ── 이미지 ────────────────────────────────────────────────────────────────
def add_image(slide, src, box, *, caption=None, frame=True, shadow=True,
              cap_h=0.24, root=None):
    """종횡비를 보존해 box 안에 이미지를 넣는다. 절대 늘리지 않는다."""
    path = Path(src)
    if not path.is_absolute() and root:
        path = Path(root) / src
    if not path.exists():
        raise SystemExit(f"이미지를 찾을 수 없음: {path}")

    x, y, cx, cy = emu_box(box)
    iw, ih = PILImage.open(path).size
    cap_emu = int(cap_h * IN) if caption else 0
    avail = cy - cap_emu
    scale = min(cx / iw, avail / ih)
    w, h = int(iw * scale), int(ih * scale)
    px, py = x + (cx - w) // 2, y + (avail - h) // 2

    pic = slide.shapes.add_picture(str(path), Emu(px), Emu(py), Emu(w), Emu(h))
    pic.name = f"이미지 {path.name}"
    sp = pic._element.spPr
    if frame:
        sp.append(parse_xml(f'<a:ln {nsdecls("a")} w="9525">'
                            f'<a:solidFill><a:srgbClr val="{CARD_LN}"/></a:solidFill></a:ln>'))
    if shadow:
        sp.append(parse_xml(f'<a:effectLst {nsdecls("a")}>'
                            f'<a:outerShdw blurRad="114300" dist="38100" dir="5400000" '
                            f'rotWithShape="0"><a:srgbClr val="1B2A41"><a:alpha val="16000"/>'
                            f"</a:srgbClr></a:outerShdw></a:effectLst>"))

    if caption:
        slide.shapes._spTree.append(S.textbox(
            name=f"캡션 {path.stem}", x=x, y=y + avail + 30000, cx=cx, cy=cap_emu,
            paras=[S.para([S.run("▲ " + caption, font=FONT, size=900, color=MUTED)],
                          align="ctr")], anchor="t"))

    dpi = round(iw / (w / IN))
    return {"name": pic.name, "px": (iw, ih), "display_in": round(w / IN, 3), "dpi": dpi}


def add_images(slide, items, box, *, gap=0.16, root=None, caption_h=0.24):
    x, y, cx, cy = box
    n = len(items)
    w = (cx - gap * (n - 1)) / n
    out = []
    for i, it in enumerate(items):
        out.append(add_image(slide, it["src"], [x + (w + gap) * i, y, w, cy],
                             caption=it.get("caption"), root=root, cap_h=caption_h))
    return out


# ── 표 ────────────────────────────────────────────────────────────────────
def _cell_border(cell, edge, color, w=9525):
    tcPr = cell._tc.get_or_add_tcPr()
    tag = f"a:ln{edge}"
    for old in tcPr.findall(qn(tag)):
        tcPr.remove(old)
    ln = parse_xml(f'<a:ln{edge} {nsdecls("a")} w="{w}" cap="flat" cmpd="sng" algn="ctr">'
                   f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln{edge}>')
    tcPr.append(ln)


def _cell_text(cell, text, *, size, color, bold=False, align="l"):
    cell.margin_left = Emu(91440)
    cell.margin_right = Emu(91440)
    cell.margin_top = Emu(45720)
    cell.margin_bottom = Emu(45720)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    pPr = p._p.get_or_add_pPr()
    pPr.set("algn", align)
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = FONT
    f.color.rgb = RGBColor.from_string(color)
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = parse_xml(f'<{tag} {nsdecls("a")} typeface="{FONT}"/>')
        rPr.append(el)


def add_table(slide, spec, box):
    x, y, cx, cy = emu_box(box)
    header, rows = spec["header"], spec["rows"]
    n_r, n_c = len(rows) + 1, len(header)
    gf = slide.shapes.add_table(n_r, n_c, Emu(x), Emu(y), Emu(cx), Emu(cy))
    gf.name = spec.get("name", "표")
    tbl = gf.table

    # 기본 테마 표 스타일(줄무늬·굵은 테두리)을 끄고 직접 칠한다
    tblPr = tbl._tbl.tblPr
    for attr, val in (("firstRow", "0"), ("bandRow", "0"), ("firstCol", "0"), ("bandCol", "0")):
        tblPr.set(attr, val)
    for sid in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(sid)

    widths = spec.get("widths") or [1] * n_c
    total = sum(widths)
    for i, col in enumerate(tbl.columns):
        col.width = Emu(int(cx * widths[i] / total))

    head_h = int(spec.get("head_h", 0.36) * IN)
    body_h = int((cy - head_h) / max(1, len(rows)))
    # 파워포인트는 글자가 들어갈 최소 높이 아래로 행을 줄이지 않는다. 그래서 box 보다
    # 표가 길어져 아래 블록을 덮는데, 도형 좌표상으로는 box 안이라 검증에 걸리지 않는다.
    # 조용히 겹치는 대신 빌드 단계에서 잡는다.
    if body_h < MIN_ROW_H:
        raise SystemExit(
            f"표 '{spec.get('name', '표')}' 행 높이 부족: "
            f"{body_h / IN:.3f}in < 최소 {MIN_ROW_H / IN:.2f}in "
            f"({len(rows)}행). box 높이를 "
            f"{(head_h + MIN_ROW_H * len(rows)) / IN:.2f}in 이상으로 늘려야 한다.")
    for i, row in enumerate(tbl.rows):
        row.height = Emu(head_h if i == 0 else body_h)

    accent = spec.get("accent", TABLE_HEAD)
    aligns = spec.get("align") or ["l"] * n_c
    for c, txt in enumerate(header):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(accent)
        _cell_text(cell, txt, size=10, color="FFFFFF", bold=True, align=aligns[c])
        for e in ("L", "R", "T", "B"):
            _cell_border(cell, e, accent)

    for r, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                "FFFFFF" if r % 2 else CARD_BG)
            _cell_text(cell, txt, size=9.5, color=BODY,
                       bold=(c == 0 and spec.get("bold_first_col", True)),
                       align=aligns[c])
            for e in ("L", "R"):
                _cell_border(cell, e, "FFFFFF", w=0)
            _cell_border(cell, "T", HAIRLINE)
            _cell_border(cell, "B", HAIRLINE)
    return gf


# ── 텍스트 ────────────────────────────────────────────────────────────────
LEVELS = {
    1: dict(size=1200, color=INK, bold=True, bullet=None, marL=0, indent=0,
            line=17, before=0),
    2: dict(size=1100, color=INK, bold=True, bullet="▪", marL=152400,
            indent=-152400, line=15.5, before=6),
    3: dict(size=1000, color=BODY, bold=False, bullet="-", marL=304800,
            indent=-152400, line=14.5, before=3),
}


def add_text(slide, items, box, *, accent=BRAND_NAVY, name="본문"):
    x, y, cx, cy = emu_box(box)
    paras = []
    for i, it in enumerate(items):
        st = LEVELS[it.get("lvl", 3)]
        paras.append(S.para(
            [S.run(it["text"], font=FONT, size=st["size"],
                   color=it.get("color", st["color"]), bold=it.get("bold", st["bold"]))],
            line_pts=st["line"], before_pts=(0 if i == 0 else st["before"]),
            marL=st["marL"], indent=st["indent"],
            bullet=st["bullet"], bullet_font="Arial",
            bullet_color=accent, bullet_size=90000))
    slide.shapes._spTree.append(
        S.textbox(name=name, x=x, y=y, cx=cx, cy=cy, paras=paras))


# ── 핵심 메시지 바 ────────────────────────────────────────────────────────
def add_note(slide, spec, box):
    x, y, cx, cy = emu_box(box)
    add = slide.shapes._spTree.append
    add(S.shape(name="핵심메시지", prst="roundRect", adj=9000, x=x, y=y, cx=cx, cy=cy,
                fill=S.grad("F2F6FC", "EAF0F8", angle=0),
                line=f'<a:ln w="9525">{S.solid("DDE5F0")}</a:ln>'))
    add(S.shape(name="핵심메시지 액센트", prst="roundRect", adj=50000,
                x=x + 100000, y=y + 110000, cx=48000, cy=cy - 220000,
                fill=S.solid(spec.get("accent", BRAND_NAVY))))
    add(S.textbox(name="핵심메시지 문구", x=x + 240000, y=y, cx=cx - 400000, cy=cy,
                  paras=[S.para([
                      S.run(spec["label"] + "     ", font=FONT, size=1100,
                            color=spec.get("accent", BRAND_NAVY), bold=True),
                      S.run(spec["text"], font=FONT, size=1050, color=INK),
                  ], line_pts=15)], anchor="ctr"))


# ── 디스패처 ──────────────────────────────────────────────────────────────
def render(slide, blocks, root=None):
    report = []
    for b in blocks:
        t = b["type"]
        if t == "image":
            report.append(add_image(slide, b["src"], b["box"], caption=b.get("caption"),
                                    frame=b.get("frame", True), root=root))
        elif t == "images":
            report += add_images(slide, b["items"], b["box"],
                                 gap=b.get("gap", 0.16), root=root)
        elif t == "table":
            add_table(slide, b, b["box"])
        elif t == "text":
            add_text(slide, b["items"], b["box"], accent=b.get("accent", BRAND_NAVY),
                     name=b.get("name", "본문"))
        elif t == "note":
            add_note(slide, b, b["box"])
        else:
            raise SystemExit(f"모르는 블록 유형: {t}")
    return report


# ── 패널 카드 (N열) ───────────────────────────────────────────────────────
def add_panels(slide, spec, box):
    x, y, cx, cy = emu_box(box)
    items = spec["items"]
    gap = int(spec.get("gap", 0.15) * IN)
    n = len(items)
    w = (cx - gap * (n - 1)) / n
    pad = 130000
    add = slide.shapes._spTree.append
    for i, it in enumerate(items):
        cx0 = x + (w + gap) * i
        color = it.get("color", BRAND_NAVY)
        dark = it.get("dark", color)
        add(S.shape(name=f"패널 {i + 1}", prst="roundRect", adj=2600,
                    x=cx0, y=y, cx=w, cy=cy, fill=S.solid(CARD_BG),
                    line=f'<a:ln w="9525">{S.solid(CARD_LN)}</a:ln>',
                    shadow=S.soft_shadow()))
        r = int(w * 0.026)
        add(S.shape(name=f"패널스트립 {i + 1}", prst="rect", x=cx0 + r, y=y,
                    cx=w - r * 2, cy=42000, fill=S.grad(color, dark, angle=0)))
        head_y = y + pad
        if it.get("no"):
            add(S.shape(name=f"패널배지 {i + 1}", prst="roundRect", adj=28000,
                        x=cx0 + pad, y=head_y, cx=340000, cy=340000,
                        fill=S.grad(color, dark, angle=2700000),
                        text_paras=[S.para([S.run(it["no"], font=FONT, size=1300,
                                                  color="FFFFFF", bold=True)], align="ctr")]))
            head_y += 430000
        add(S.textbox(name=f"패널제목 {i + 1}", x=cx0 + pad, y=head_y,
                      cx=w - pad * 2, cy=460000,
                      paras=[S.para([S.run(it["title"], font=FONT, size=1250,
                                           color=INK, bold=True)], line_pts=16)]))
        ty = head_y + 400000
        if it.get("subtitle"):
            add(S.textbox(name=f"패널부제 {i + 1}", x=cx0 + pad, y=ty,
                          cx=w - pad * 2, cy=250000,
                          paras=[S.para([S.run(it["subtitle"], font=FONT, size=950,
                                               color=color, bold=True)], line_pts=13)]))
            ty += 250000
        add(S.rule(name=f"패널선 {i + 1}", x=cx0 + pad, y=ty, cx=w - pad * 2, color=HAIRLINE))
        add(S.textbox(
            name=f"패널본문 {i + 1}", x=cx0 + pad, y=ty + 80000,
            cx=w - pad * 2, cy=y + cy - (ty + 80000) - 110000,
            paras=[S.para([S.run(t, font=FONT, size=1000, color=BODY)],
                          line_pts=13.5, before_pts=(0 if j == 0 else 5),
                          marL=139700, indent=-139700,
                          bullet="▪", bullet_font="Arial", bullet_color=color,
                          bullet_size=90000)
                   for j, t in enumerate(it.get("items", []))]))


# ── 흐름도 (가로 단계 + 화살표) ───────────────────────────────────────────
def add_flow(slide, spec, box):
    x, y, cx, cy = emu_box(box)
    nodes = spec["nodes"]
    n = len(nodes)
    gap = int(spec.get("gap", 0.42) * IN)
    node_h = int(cy * spec.get("node_ratio", 0.60))
    w = (cx - gap * (n - 1)) / n
    add = slide.shapes._spTree.append
    for i, nd in enumerate(nodes):
        nx = x + (w + gap) * i
        color = nd.get("color", BRAND_NAVY)
        dark = nd.get("dark", color)
        paras = [S.para([S.run(nd["label"], font=FONT, size=1200, color="FFFFFF",
                               bold=True)], align="ctr", line_pts=15)]
        if nd.get("sub"):
            paras.append(S.para([S.run(nd["sub"], font=FONT, size=900, color="E8EEF7")],
                                align="ctr", line_pts=12, before_pts=2))
        add(S.shape(name=f"흐름 {i + 1}", prst="roundRect", adj=9000,
                    x=nx, y=y, cx=w, cy=node_h,
                    fill=S.grad(color, dark, angle=5400000),
                    line=f'<a:ln w="9525">{S.solid(dark)}</a:ln>',
                    text_paras=paras, ins=(45720, 45720, 45720, 45720),
                    shadow=S.soft_shadow()))
        if i < n - 1:
            ax = nx + w + (gap - 150000) / 2
            # 화살표 색은 기본 회색. 단계 사이가 '단절'인 경우 노드에서 덮어쓴다.
            add(S.shape(name=f"화살표 {i + 1}", prst="triangle", rot=5400000,
                        x=ax, y=y + node_h / 2 - 90000, cx=150000, cy=180000,
                        fill=S.solid(nd.get("arrow_color", "A9B4C6"))))
            if nd.get("arrow_label"):
                add(S.textbox(name=f"화살표라벨 {i + 1}",
                              x=nx + w - 500000, y=y + node_h + 60000,
                              cx=gap + 1000000, cy=260000,
                              paras=[S.para([S.run(nd["arrow_label"], font=FONT,
                                                   size=900, color="6B7688")], align="ctr")]))


def _render_one(slide, b, root):
    t = b["type"]
    if t == "image":
        return [add_image(slide, b["src"], b["box"], caption=b.get("caption"),
                          frame=b.get("frame", True), root=root)]
    if t == "images":
        return add_images(slide, b["items"], b["box"], gap=b.get("gap", 0.16), root=root)
    if t == "table":
        add_table(slide, b, b["box"])
    elif t == "text":
        add_text(slide, b["items"], b["box"], accent=b.get("accent", BRAND_NAVY),
                 name=b.get("name", "본문"))
    elif t == "note":
        add_note(slide, b, b["box"])
    elif t == "panels":
        add_panels(slide, b, b["box"])
    elif t == "flow":
        add_flow(slide, b, b["box"])
    else:
        # 보고자료용 확장 블록(colhead·vsteps·cycle·strip)에 위임한다
        import blocks_report as R
        if t not in R.RENDERERS:
            raise SystemExit(f"모르는 블록 유형: {t}")
        R.render_block(slide, b)
    return []


def render(slide, blocks, root=None):   # noqa: F811  (위 정의를 대체)
    out = []
    for b in blocks:
        out += _render_one(slide, b, root)
    return out
