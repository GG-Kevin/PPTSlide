#!/usr/bin/env python3
"""회사 표준 포맷 위에 로드맵형 슬라이드를 생성한다.

  python3 tools/build_slide.py content/intro_dev_schedule.json "output/이름.pptx"

레이아웃(위 → 아래)
  A. 제목 + ❖ 대분류 (템플릿 플레이스홀더 — 포맷 상속)
  B. 로드맵 간트   : 월 축 + 단계별 기간 바 (기간·공백기가 눈에 보인다)
  C. 단계 카드 5장 : 번호 배지 · 기간 칩 · 구분선 · 세부 항목
  D. 핵심 방향 바  : 한 줄 요약 (회사 남색 #014099 액센트)

포맷 규칙
  * templates/company_format.pptx 를 열어 레이아웃을 상속받는다.
  * 상단 모티프 바 / 우하단 페이지번호 / 제목 스타일은 건드리지 않는다.
  * 본문 글꼴은 회사 폰트가 없는 PC를 고려해 테마 폰트인 '맑은 고딕'을 쓴다.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

sys.path.insert(0, str(Path(__file__).resolve().parent))
import shapes as S  # noqa: E402
import blocks as B  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "company_format.pptx"
SPEC = json.loads((ROOT / "format" / "format_spec.json").read_text(encoding="utf-8"))

# ── 글꼴 ───────────────────────────────────────────────────────────────────
# 회사 폰트 '현대하모니 L' 이 없는 PC 에서도 깨지지 않도록 테마 폰트를 쓴다.
FONT = "맑은 고딕"

# ── 색 (기준 이미지의 보라→청록 그라데이션에서 파생) ────────────────────────
INK = "1B2A41"        # 제목/강조 텍스트
BODY = "44536A"       # 본문 텍스트
MUTED = "8A94A6"      # 보조 텍스트
HAIRLINE = "E3E9F0"   # 구분선
GRID = "EFF2F7"       # 간트 격자
CARD_BG = "F7F9FC"    # 카드 배경
CARD_LN = "E4EAF2"    # 카드 테두리
BRAND_NAVY = "014099"  # 회사 모티프 바 남색

# ── 세로 배치 (EMU) — 본문 플레이스홀더 764704 ~ 6453188 안에서 끝난다 ──────
CONTENT_L = 426800          # ❖ 불릿과 같은 좌측 기준선
CONTENT_R = 11856640        # 본문 플레이스홀더 우측 끝
CONTENT_W = CONTENT_R - CONTENT_L

META_Y = 820000

BAND_Y = 1240000            # 간트 시작
AXIS_LABEL_H = 250000
BAND_RULE_Y = 1510000
ROW_Y0 = 1570000
ROW_PITCH = 395000
BAR_H = 185000
LABEL_COL_W = 2350000
AXIS_GAP = 180000
ROW_BAND = "FBFCFE"   # 간트 짝수행 배경

CARD_Y = 3760000
CARD_H = 1960000
CARD_GAP = 137160

KEY_Y = 5880000
KEY_H = 520000

PAD = 130000


def sp(shape):
    return shape


def build_roadmap(slide, c):
    add = slide.shapes._spTree.append
    S.reset_ids()

    # ── A. 우측 메타 라인 ────────────────────────────────────────────────
    add(S.textbox(
        name="메타", x=CONTENT_R - 6000000, y=META_Y, cx=6000000, cy=280000,
        paras=[S.para([S.run(c["meta"], font=FONT, size=1000, color=MUTED)], align="r")],
        anchor="ctr"))

    # ── B. 로드맵 간트 ───────────────────────────────────────────────────
    axis_x = CONTENT_L + LABEL_COL_W + AXIS_GAP
    axis_w = CONTENT_R - axis_x
    n = len(c["axis"])
    mw = axis_w / n
    rows_bottom = ROW_Y0 + ROW_PITCH * len(c["phases"])

    # 월 라벨
    for i, m in enumerate(c["axis"]):
        add(S.textbox(name=f"월 {m}", x=axis_x + mw * i, y=BAND_Y, cx=mw, cy=AXIS_LABEL_H,
                      paras=[S.para([S.run(m, font=FONT, size=950, color=MUTED)], align="ctr")],
                      anchor="ctr"))
    # 축 구분선
    add(S.rule(name="축 구분선", x=CONTENT_L, y=BAND_RULE_Y, cx=CONTENT_W, color=HAIRLINE))
    # 월 격자
    for i in range(n + 1):
        add(S.shape(name=f"격자 {i}", prst="rect", x=axis_x + mw * i, y=BAND_RULE_Y,
                    cx=9525, cy=rows_bottom - BAND_RULE_Y, fill=S.solid(GRID)))

    for r, ph in enumerate(c["phases"]):
        row_y = ROW_Y0 + ROW_PITCH * r
        if r % 2 == 1:   # 짝수행 옅은 배경 — 행 추적이 쉬워진다
            add(S.shape(name=f"행배경 {r}", prst="rect", x=CONTENT_L, y=row_y,
                        cx=CONTENT_W, cy=ROW_PITCH, fill=S.solid(ROW_BAND)))
        # 좌측 라벨: 번호(단계색) + 이름
        add(S.textbox(
            name=f"간트라벨 {ph['no']}", x=CONTENT_L, y=row_y, cx=LABEL_COL_W, cy=ROW_PITCH,
            paras=[S.para([
                S.run(ph["no"] + "   ", font=FONT, size=1050, color=ph["color"], bold=True),
                S.run(ph["name"], font=FONT, size=1050, color=BODY),
            ])],
            anchor="ctr"))
        # 기간 바 (여러 구간 지원 — 4단계처럼 중간에 공백기가 있는 경우)
        bar_y = row_y + (ROW_PITCH - BAR_H) / 2
        prev_end = None
        for s0, s1 in ph["spans"]:
            bx, bw = axis_x + mw * s0, mw * (s1 - s0)
            if prev_end is not None:      # 공백기 연결선(점선)
                add(S.shape(name=f"공백기 {ph['no']}", prst="line",
                            x=prev_end, y=bar_y + BAR_H / 2, cx=bx - prev_end, cy=0,
                            line=f'<a:ln w="9525"><a:solidFill><a:srgbClr val="{ph["color"]}">'
                                 f'<a:alpha val="45000"/></a:srgbClr></a:solidFill>'
                                 f'<a:prstDash val="sysDot"/></a:ln>'))
            add(S.shape(name=f"기간바 {ph['no']}", prst="roundRect", adj=50000,
                        x=bx, y=bar_y, cx=bw, cy=BAR_H,
                        fill=S.grad(ph["color"], ph["dark"], angle=0)))
            prev_end = bx + bw

    # ── C. 단계 카드 ─────────────────────────────────────────────────────
    k = len(c["phases"])
    card_w = (CONTENT_W - CARD_GAP * (k - 1)) / k
    for i, ph in enumerate(c["phases"]):
        cx0 = CONTENT_L + (card_w + CARD_GAP) * i
        add(S.shape(name=f"카드 {ph['no']}", prst="roundRect", adj=2600,
                    x=cx0, y=CARD_Y, cx=card_w, cy=CARD_H,
                    fill=S.solid(CARD_BG),
                    line=f'<a:ln w="9525">{S.solid(CARD_LN)}</a:ln>',
                    shadow=S.soft_shadow()))
        # 카드 상단 색 스트립 (간트 바와 같은 색 → 시각적 연결)
        r_emu = int(card_w * 0.026)
        add(S.shape(name=f"카드스트립 {ph['no']}", prst="rect",
                    x=cx0 + r_emu, y=CARD_Y, cx=card_w - r_emu * 2, cy=42000,
                    fill=S.grad(ph["color"], ph["dark"], angle=0)))
        # 번호 배지
        add(S.shape(name=f"배지 {ph['no']}", prst="roundRect", adj=28000,
                    x=cx0 + PAD, y=CARD_Y + PAD, cx=340000, cy=340000,
                    fill=S.grad(ph["color"], ph["dark"], angle=2700000),
                    text_paras=[S.para([S.run(ph["no"], font=FONT, size=1300,
                                              color="FFFFFF", bold=True)], align="ctr")],
                    ins=(0, 0, 0, 0)))
        # 기간 칩
        chip_w = 1120000
        add(S.shape(name=f"기간칩 {ph['no']}", prst="roundRect", adj=45000,
                    x=cx0 + card_w - PAD - chip_w, y=CARD_Y + PAD + 55000,
                    cx=chip_w, cy=230000,
                    fill=S.solid(ph["tint"]),
                    text_paras=[S.para([S.run(ph["period"], font=FONT, size=850,
                                              color=ph["dark"], bold=True)], align="ctr")],
                    ins=(0, 0, 0, 0)))
        # 단계명
        add(S.textbox(
            name=f"카드제목 {ph['no']}", x=cx0 + PAD, y=CARD_Y + 545000,
            cx=card_w - PAD * 2, cy=500000,
            paras=[S.para([S.run(ph["name"], font=FONT, size=1300, color=INK, bold=True)],
                          line_pts=17)]))
        # 구분선
        add(S.rule(name=f"카드선 {ph['no']}", x=cx0 + PAD, y=CARD_Y + 1090000,
                   cx=card_w - PAD * 2, color=HAIRLINE))
        # 세부 항목. 상단 정렬을 유지해 카드 5장의 첫 줄이 가로로 맞게 한다
        # (하단 여백 편차는 카드 높이를 항목 3개에 맞춰 줄이는 쪽으로 해결)
        add(S.textbox(
            name=f"카드본문 {ph['no']}", x=cx0 + PAD, y=CARD_Y + 1160000,
            cx=card_w - PAD * 2, cy=CARD_H - 1160000 - 120000,
            paras=[S.para([S.run(t, font=FONT, size=1000, color=BODY)],
                          line_pts=12.5, before_pts=(0 if j == 0 else 4.5),
                          marL=139700, indent=-139700,
                          bullet="▪", bullet_font="Arial", bullet_color=ph["color"],
                          bullet_size=90000)
                   for j, t in enumerate(ph["points"])]))

    # ── D. 핵심 방향 바 ──────────────────────────────────────────────────
    km = c["key_message"]
    add(S.shape(name="핵심방향", prst="roundRect", adj=9000,
                x=CONTENT_L, y=KEY_Y, cx=CONTENT_W, cy=KEY_H,
                fill=S.grad("F2F6FC", "EAF0F8", angle=0),
                line=f'<a:ln w="9525">{S.solid("DDE5F0")}</a:ln>'))
    add(S.shape(name="핵심방향 액센트", prst="roundRect", adj=50000,
                x=CONTENT_L + 100000, y=KEY_Y + 110000, cx=48000, cy=KEY_H - 220000,
                fill=S.solid(BRAND_NAVY)))
    add(S.textbox(
        name="핵심방향 문구", x=CONTENT_L + 240000, y=KEY_Y,
        cx=CONTENT_W - 400000, cy=KEY_H,
        paras=[S.para([
            S.run(km["label"] + "     ", font=FONT, size=1100, color=BRAND_NAVY, bold=True),
            S.run(km["text"], font=FONT, size=1050, color=INK),
        ], line_pts=15)],
        anchor="ctr"))


def delete_all_slides(prs):
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.get(qn("r:id")))
        lst.remove(sld)


def get_layout(prs, name):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise SystemExit(f"레이아웃 '{name}' 을 포맷에서 찾을 수 없음")


def build(content_path, out_path):
    c = json.loads(Path(content_path).read_text(encoding="utf-8"))
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)
    slide = prs.slides.add_slide(get_layout(prs, c.get("layout", SPEC["body_slide_layout"])))

    # 제목 · ❖ 대분류는 템플릿 플레이스홀더에 그대로 넣어 포맷을 상속받는다
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = c["title"]
    body = next(s for s in slide.placeholders
                if s.element.find(f".//{qn('p:ph')}") is not None
                and s.element.find(f".//{qn('p:ph')}").get("idx") == "10")
    body.text_frame.paragraphs[0].add_run().text = c["heading"]

    design = c.get("design", "roadmap")
    if design == "roadmap":
        build_roadmap(slide, c)
    elif design == "blocks":
        # 초안을 블록 목록으로 옮긴 슬라이드 (이미지 · 표 · 텍스트 자유 조합)
        if c.get("meta"):
            slide.shapes._spTree.append(S.textbox(
                name="메타", x=CONTENT_R - 6000000, y=META_Y, cx=6000000, cy=280000,
                paras=[S.para([S.run(c["meta"], font=FONT, size=1000, color=MUTED)],
                              align="r")], anchor="ctr"))
        for item in B.render(slide, c["blocks"], root=ROOT):
            print("  이미지:", item["name"], f"{item['px'][0]}x{item['px'][1]}px",
                  f"표시 {item['display_in']}in", f"{item['dpi']}dpi")
    else:
        raise SystemExit(f"모르는 design: {design}")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"wrote {out_path}")


if __name__ == "__main__":
    build(sys.argv[1], sys.argv[2])
