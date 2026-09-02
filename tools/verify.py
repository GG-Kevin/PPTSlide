#!/usr/bin/env python3
"""생성된 pptx 가 회사 표준 포맷을 지켰는지 기계적으로 검증한다.

  python3 tools/verify.py output/Introduction_개발일정.pptx

A. 구조 검증 : pptx XML 을 format/format_spec.json 과 대조
B. 화면 검증 : 렌더링 이미지를 format/reference_intro_slide.png 와 대조
    (모티프 바 위치·색, 흰 배경, 페이지번호, 제목 위치, 좌우 여백)

종료코드 0 = 합격, 1 = 불합격. 리포트는 output/verify_report.json 에도 남는다.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SPEC = json.loads((ROOT / "format" / "format_spec.json").read_text(encoding="utf-8"))
REFERENCE = ROOT / "format" / "reference_intro_slide.png"
# 회사 폰트 + 테마 폰트(맑은 고딕, 회사 폰트가 없는 PC 대비) + 불릿용 심볼 폰트
ALLOWED_FONTS = {"현대하모니 L", "맑은 고딕", "Arial", "Wingdings", None}

# 본문 콘텐츠가 넘어서면 안 되는 경계 = 본문 플레이스홀더 bbox (포맷에서 계산)
_BODY_PH = next(s for s in SPEC["layouts"][SPEC["body_slide_layout"]]["shapes"]
                if s["ph"] and s["ph"].get("idx") == "10")
BODY_BOX = (_BODY_PH["left"], _BODY_PH["top"],
            _BODY_PH["left"] + _BODY_PH["width"], _BODY_PH["top"] + _BODY_PH["height"])

results = []


def check(name, ok, detail=""):
    results.append({"check": name, "ok": bool(ok), "detail": detail})
    return ok


# ─────────────────────────── A. 구조 검증 ───────────────────────────
def verify_structure(pptx_path):
    prs = Presentation(str(pptx_path))
    ss = SPEC["slide_size_emu"]
    check("슬라이드 크기", prs.slide_width == ss["cx"] and prs.slide_height == ss["cy"],
          f"{prs.slide_width}x{prs.slide_height} (기준 {ss['cx']}x{ss['cy']})")

    check("슬라이드 1장", len(prs.slides.__iter__.__self__._sldIdLst) >= 1,
          f"{len(prs.slides._sldIdLst)}장")

    master_names = [m.name for m in prs.slide_masters]
    check("마스터 유지", len(prs.slide_masters) == 1, str(master_names))

    layout_names = {l.name for l in prs.slide_masters[0].slide_layouts}
    check("레이아웃 세트 유지", layout_names == set(SPEC["layouts"].keys()),
          f"누락 {set(SPEC['layouts']) - layout_names} / 추가 {layout_names - set(SPEC['layouts'])}")

    for i, slide in enumerate(prs.slides, 1):
        p = f"슬라이드{i}"
        check(f"{p} 본문 레이아웃 사용", slide.slide_layout.name == SPEC["body_slide_layout"],
              slide.slide_layout.name)

        phs = {}
        for sh in slide.shapes:
            ph = sh.element.find(f".//{qn('p:ph')}")
            if ph is not None:
                phs[ph.get("type") or f"idx{ph.get('idx')}"] = sh
        check(f"{p} 제목 플레이스홀더 사용", "title" in phs, str(list(phs)))
        check(f"{p} 본문 플레이스홀더 사용", any(k.startswith("idx10") or k == "idx10" for k in phs)
              or any((sh.element.find(f".//{qn('p:ph')}") is not None
                      and sh.element.find(f".//{qn('p:ph')}").get("idx") == "10")
                     for sh in slide.shapes), str(list(phs)))

        # 제목 텍스트에 로컬 폰트/크기/색 덮어쓰기가 없어야 한다
        title = prs.slides[i - 1].shapes.title
        overrides = []
        for r in title.element.iter(qn("a:rPr")):
            for attr in ("sz", "b", "i"):
                if r.get(attr) is not None:
                    overrides.append(f"{attr}={r.get(attr)}")
            if r.find(qn("a:latin")) is not None:
                overrides.append("latin")
            if r.find(qn("a:solidFill")) is not None:
                overrides.append("fill")
        check(f"{p} 제목 스타일 상속(로컬 덮어쓰기 없음)", not overrides, ", ".join(overrides))

        # 폰트 화이트리스트
        bad_fonts = set()
        for latin in slide.shapes._spTree.iter(qn("a:latin")):
            if latin.get("typeface") not in ALLOWED_FONTS:
                bad_fonts.add(latin.get("typeface"))
        for ea in slide.shapes._spTree.iter(qn("a:ea")):
            if ea.get("typeface") not in ALLOWED_FONTS:
                bad_fonts.add(ea.get("typeface"))
        check(f"{p} 승인 폰트만 사용", not bad_fonts, f"허용외 {sorted(bad_fonts)}")

        # 모티프 바 / 페이지번호를 슬라이드에서 다시 그리지 않았는지
        motif_top, motif_h = SPEC["motif_bar"]["top_emu"], SPEC["motif_bar"]["height_emu"]
        clashes = []
        for sh in slide.shapes:
            if sh.top is None:
                continue
            if sh.top < motif_top + motif_h and sh.top + (sh.height or 0) > motif_top \
               and (sh.width or 0) > prs.slide_width * 0.5 and sh.element.find(f".//{qn('p:ph')}") is None:
                clashes.append(sh.name)
            if sh.top > prs.slide_height * 0.93 and sh.left > prs.slide_width * 0.88:
                clashes.append(sh.name)
        check(f"{p} 모티프바/페이지번호 영역 침범 없음", not clashes, ", ".join(clashes))

        # 직접 추가한 도형이 본문 영역(플레이스홀더 bbox)을 벗어나지 않았는지
        bx0, by0, bx1, by1 = BODY_BOX
        tol = 20000
        strays = [sh.name for sh in slide.shapes
                  if sh.element.find(f".//{qn('p:ph')}") is None and sh.left is not None
                  and (sh.left < bx0 - tol or sh.top < by0 - tol
                       or sh.left + sh.width > bx1 + tol
                       or sh.top + sh.height > by1 + tol)]
        check(f"{p} 도형이 본문 영역 안에 있음", not strays, ", ".join(strays[:5]))

        # 카드끼리 겹치지 않는지 (좌우 배치가 무너지면 바로 잡힌다)
        cards = sorted([sh for sh in slide.shapes if sh.name.startswith("카드 ")],
                       key=lambda s: s.left)
        laps = [f"{a.name}/{b.name}" for a, b in zip(cards, cards[1:])
                if a.left + a.width > b.left + 1000]
        check(f"{p} 카드 겹침 없음", not laps, f"{len(cards)}장, " + ", ".join(laps))

        # 간트 기간 바가 월 축 범위 안에 있는지
        grids = [sh for sh in slide.shapes if sh.name.startswith("격자 ")]
        bars = [sh for sh in slide.shapes if sh.name.startswith("기간바 ")]
        if grids and bars:
            ax0 = min(g.left for g in grids)
            ax1 = max(g.left + g.width for g in grids)
            over = [b.name for b in bars if b.left < ax0 - 1000 or b.left + b.width > ax1 + 1000]
            check(f"{p} 기간 바가 월 축 범위 내", not over, f"{len(bars)}개, " + ", ".join(over))

        # 슬라이드 밖으로 나간 도형 없음
        out = [sh.name for sh in slide.shapes
               if sh.left is not None and (sh.left < 0 or sh.top < 0
                                           or sh.left + sh.width > prs.slide_width
                                           or sh.top + sh.height > prs.slide_height)]
        check(f"{p} 슬라이드 경계 이탈 없음", not out, ", ".join(out))

        # 본문 내용이 제목/모티프 아래에서 시작
        body_tops = [sh.top for sh in slide.shapes
                     if sh.element.find(f".//{qn('p:ph')}") is None and sh.top is not None]
        if body_tops:
            check(f"{p} 콘텐츠가 모티프 바 아래", min(body_tops) > motif_top + motif_h,
                  f"최상단 {min(body_tops)} EMU")


# ─────────────────────────── B. 화면 검증 ───────────────────────────
def motif_profile(img):
    from collections import Counter
    W, H = img.size
    best = None
    for y in range(int(H * 0.04), int(H * 0.15)):
        row = [img.getpixel((x, y)) for x in range(2, W - 2)]
        c = Counter(row)
        colored = [(p, n) for p, n in c.items() if not (p[0] > 240 and p[1] > 240 and p[2] > 240)]
        colored.sort(key=lambda t: -t[1])
        if len(colored) >= 2 and colored[0][1] + colored[1][1] > (W - 4) * 0.9:
            best = (y, colored[0][0], colored[1][0], colored[0][1], colored[1][1])
            break
    if not best:
        return None
    y, c1, c2, n1, n2 = best
    row = [img.getpixel((x, y)) for x in range(2, W - 2)]
    split = next((i for i, p in enumerate(row) if p == c2), 0) / len(row)
    return {"y_ratio": y / H, "left_rgb": list(c1), "right_rgb": list(c2), "split_ratio": round(split, 3)}


def verify_render(png_path):
    from PIL import Image

    img = Image.open(png_path).convert("RGB")
    ref = Image.open(REFERENCE).convert("RGB")
    W, H = img.size

    # 본문 영역 아래 여백(콘텐츠도 페이지번호도 없는 곳)에서 배경색을 본다
    bg_pts = [(0.40, 0.97), (0.012, 0.50), (0.60, 0.985)]
    bg = [img.getpixel((int(W * fx), int(H * fy))) for fx, fy in bg_pts]
    check("배경 흰색", all(p[0] > 245 and p[1] > 245 and p[2] > 245 for p in bg), str(bg))

    got, want = motif_profile(img), motif_profile(ref)
    if check("상단 모티프 바 검출", got is not None and want is not None, f"{got}"):
        check("모티프 바 세로위치", abs(got["y_ratio"] - want["y_ratio"]) < 0.01,
              f"{got['y_ratio']:.4f} vs 기준 {want['y_ratio']:.4f}")
        check("모티프 바 좌측 남색", max(abs(a - b) for a, b in zip(got["left_rgb"], want["left_rgb"])) <= 6,
              f"{got['left_rgb']} vs {want['left_rgb']}")
        check("모티프 바 우측 녹색", max(abs(a - b) for a, b in zip(got["right_rgb"], want["right_rgb"])) <= 6,
              f"{got['right_rgb']} vs {want['right_rgb']}")
        check("모티프 바 색 분할점", abs(got["split_ratio"] - want["split_ratio"]) < 0.02,
              f"{got['split_ratio']} vs {want['split_ratio']}")

    def dark_bbox(im, x0, y0, x1, y1, skip=None):
        """어두운 픽셀의 bbox 를 비율로 반환. skip 은 제외할 사각형(비율)."""
        xs, ys = [], []
        iw, ih = im.size
        for y in range(int(y0), int(y1)):
            for x in range(int(x0), int(x1)):
                if skip and skip[0] <= x / iw <= skip[2] and skip[1] <= y / ih <= skip[3]:
                    continue
                if sum(im.getpixel((x, y))) < 400:
                    xs.append(x); ys.append(y)
        return (min(xs) / iw, min(ys) / ih, max(xs) / iw, max(ys) / ih) if xs else None

    t_got = dark_bbox(img, 5, H * 0.01, W * 0.6, H * 0.075)
    t_ref = dark_bbox(ref, 5, ref.size[1] * 0.01, ref.size[0] * 0.6, ref.size[1] * 0.075)
    if check("제목 텍스트 존재", t_got is not None, str(t_got)):
        check("제목 좌측 정렬 위치", abs(t_got[0] - t_ref[0]) < 0.012,
              f"x={t_got[0]:.4f} vs 기준 {t_ref[0]:.4f}")
        check("제목 세로 위치", abs(t_got[1] - t_ref[1]) < 0.015,
              f"y={t_got[1]:.4f} vs 기준 {t_ref[1]:.4f}")

    pn = dark_bbox(img, W * 0.95, H * 0.94, W - 2, H - 2)
    check("우하단 페이지 번호", pn is not None, str(pn))

    # 본문 허용 경계는 포맷의 본문 플레이스홀더에서 계산한다 (임의의 상수 금지)
    right_max = BODY_BOX[2] / SPEC["slide_size_emu"]["cx"] + 0.005
    bottom_max = BODY_BOX[3] / SPEC["slide_size_emu"]["cy"] + 0.005
    page_num = (0.95, 0.94, 1.0, 1.0)   # 우하단 페이지 번호는 본문이 아니므로 제외
    body = dark_bbox(img, 2, H * 0.10, W - 2, H * 0.99, skip=page_num)
    ref_body = dark_bbox(ref, 2, ref.size[1] * 0.10, ref.size[0] - 2,
                         ref.size[1] * 0.99, skip=page_num)
    if check("본문 콘텐츠 존재", body is not None, str(body)):
        check("본문 좌측 여백", abs(body[0] - ref_body[0]) < 0.012,
              f"x={body[0]:.4f} vs 기준 {ref_body[0]:.4f}")
        check("본문 우측 여백 초과 없음", body[2] <= right_max,
              f"right={body[2]:.4f} (한계 {right_max:.4f})")
        check("본문 하단 넘침 없음", body[3] <= bottom_max,
              f"bottom={body[3]:.4f} (한계 {bottom_max:.4f})")


def main():
    pptx_path = Path(sys.argv[1])
    verify_structure(pptx_path)
    from render import render
    pngs = render(pptx_path, ROOT / "output" / "preview", dpi=102)
    verify_render(pngs[0])

    failed = [r for r in results if not r["ok"]]
    report = {"target": str(pptx_path), "passed": len(results) - len(failed),
              "failed": len(failed), "results": results}
    (ROOT / "output" / "verify_report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    for r in results:
        print(f"[{'PASS' if r['ok'] else 'FAIL'}] {r['check']}" + (f"  — {r['detail']}" if r["detail"] else ""))
    print(f"\n{len(results) - len(failed)}/{len(results)} 통과")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    sys.exit(main())
